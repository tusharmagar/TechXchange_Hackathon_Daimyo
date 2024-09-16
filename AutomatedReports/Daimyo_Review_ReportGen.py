from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx import Presentation
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import json
import win32com.client
import os
import matplotlib.dates as mdates
from datetime import datetime, timedelta
import pytz
from sqlalchemy import create_engine, text
from sqlalchemy.orm import sessionmaker
from azure.storage.blob import BlobServiceClient
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from twilio.rest import Client
from langchain_ibm import ChatWatsonx
import json


# Set up environment variables for Langchain
os.environ["LANGCHAIN_TRACING_V2"] = "true"
os.environ["LANGCHAIN_PROJECT"] = "<YOUR_LANGCHAIN_PROJECT>"
os.environ["LANGCHAIN_ENDPOINT"] = "<YOUR_LANGCHAIN_ENDPOINT>"
os.environ["LANGCHAIN_API_KEY"] = "<YOUR_LANGCHAIN_API_KEY>"

# Initialize the Granite model
parameters = {
    "decoding_method": "greedy",
    "max_new_tokens": 900,
    "repetition_penalty": 1,
    "min_new_tokens":1
}
Watson_llm = ChatWatsonx(
    model_id="ibm/granite-34b-code-instruct",
    url="https://us-south.ml.cloud.ibm.com",
    apikey="YOUR API KEY",
    project_id="YOUR PROJECT ID",
    params=parameters,
)
summary_prompt = ChatPromptTemplate.from_template(
    """
        You are a Review Summarization assistant. Your task is to provide a concise and coherent summary of the reviews provided.
        INSTRUCTIONS:
        1. Always use ONLY the information from the reviews to generate summaries.
        2. Summary has a maximum of 200 words.
        3. Focuses on key points, sentiment, and common themes in the reviews.

        THIS BELOW STEP IS VERY IMPORTANT
        If no reviews are given, always respond with: "No reviews available."

        Reviews: {Reviews}
        Summary:
    """
)
summary_chain = summary_prompt | Watson_llm | StrOutputParser()

# Database connection using SQLAlchemy
engine = create_engine('mssql+pymssql://<DB_USER>:<DB_PASSWORD>@<DB_HOST>/<DB_NAME>', echo=False)

# Azure Blob Storage setup
blob_service_client = BlobServiceClient.from_connection_string("<AZURE_CONNECTION_STRING>")
container_name = "<CONTAINER_NAME>"

# Twilio configuration
account_sid = '<TWILIO_ACCOUNT_SID>'
auth_token = '<TWILIO_AUTH_TOKEN>'
twilio_client = Client(account_sid, auth_token)

# Messaging service SID
messaging_service_sid = '<TWILIO_MESSAGING_SERVICE_SID>'

# WhatsApp template content SID
content_sid = "<WHATSAPP_CONTENT_SID>"

# Map of restaurant names to WhatsApp phone numbers
PHONE_NUMBERS = {
    'Restaurant 1': ['whatsapp:<PHONE_NUMBER_1>', 'whatsapp:<PHONE_NUMBER_2>'],
    'Restaurant 2': ['whatsapp:<PHONE_NUMBER_1>', 'whatsapp:<PHONE_NUMBER_2>']
}

# Function to fetch data from the database
def fetch_data_from_db():
    query = text('''
        SELECT 
            place_name, 
            google_id, 
            review_id, 
            author_title, 
            author_id, 
            review_text, 
            review_questions, 
            review_link, 
            review_rating, 
            review_datetime_utc, 
            sentiment_score
        FROM GoogleReviews
    ''')
    with engine.connect() as connection:
        data = pd.read_sql(query, connection)
    return data

# Function to fetch restaurant names from the database
def fetch_restro_from_db():
    query = text('''
        SELECT DISTINCT place_name 
        FROM GoogleReviews
    ''')
    with engine.connect() as connection:
        result = connection.execute(query)
        restaurant_names = [row[0] for row in result]  # Extract the place_name values
    return restaurant_names

def filter_by_restaurant(data, restaurant_name):
    return data[data['place_name'].str.contains(restaurant_name, case=False, na=False)]

colors = {
    "first_row": RGBColor(112, 173, 71),  # #70AD47
    "even_row": RGBColor(213, 227, 207),  # #D5E3CF
    "odd_row": RGBColor(235, 241, 233)    # #EBF1E9
    }

# Function to safely parse the JSON in 'review_questions'
def parse_review_questions(review_questions):
    try:
        return json.loads(review_questions.replace("null", "None"))
    except (json.JSONDecodeError, TypeError):
        return {}


# Function to check if the report has already been generated this week
def check_report_generated(restaurant_name, report_type):
    query = text('''
        SELECT last_report_date 
        FROM SentReport 
        WHERE place_name = :restaurant_name AND report_type = :report_type
    ''')
    with engine.connect() as connection:
        result = connection.execute(query, {'restaurant_name': restaurant_name, 'report_type': report_type}).fetchone()
    if result is None:
        return False  # No report has ever been generated for this restaurant and report type
    last_report_date = result[0]
    
    # Check if the report was generated this week
    current_week_start = datetime.now() - timedelta(days=datetime.now().weekday())  # Start of the current week
    print(f"Last report date for {restaurant_name}: {last_report_date}") #-_-
    return last_report_date >= current_week_start.date()

# Function to update the last report generation date
def update_report_generated_date(restaurant_name, report_type):
    query = text('''
        MERGE INTO SentReport AS target
        USING (SELECT :restaurant_name AS place_name, :report_type AS report_type, :current_date AS last_report_date) AS source
        ON target.place_name = source.place_name AND target.report_type = source.report_type
        WHEN MATCHED THEN UPDATE SET target.last_report_date = source.last_report_date
        WHEN NOT MATCHED THEN INSERT (place_name, report_type, last_report_date) VALUES (source.place_name, source.report_type, source.last_report_date);
    ''')

    print(f"Query: {query}")
    print(f"Parameters: {restaurant_name}, {report_type}, {datetime.now().date()}")

    Session = sessionmaker(bind=engine)
    session = Session()

    try:
        session.execute(query, {
            'restaurant_name': restaurant_name, 
            'report_type': report_type, 
            'current_date': datetime.now().date()
        })
        session.commit()
        print(f"Updated report date for {restaurant_name}.")
    except Exception as e:
        session.rollback()
        print(f"Error updating report date for {restaurant_name}: {e}")
    finally:
        session.close()




def send_report_via_whatsapp(restaurant_name, pdf_blob_url):
    # Log the full URL first
    print(f"Full PDF Blob URL: {pdf_blob_url}")

    # Extract the filename from the URL
    pdf_blob_filename = pdf_blob_url.split('/')[-1]
    
    # Log the extracted filename to verify correctness
    print(f"Extracted PDF Blob Filename: {pdf_blob_filename}")

    # Check if the filename contains spaces and replace with '%20'
    if ' ' in pdf_blob_filename:
        pdf_blob_filename = pdf_blob_filename.replace(' ', '%20')
        print(f"Filename after replacing spaces: {pdf_blob_filename}")

    # Get the phone numbers for the restaurant
    phone_numbers = PHONE_NUMBERS.get(restaurant_name, [])
    
    # Log the restaurant name and the phone numbers
    print(f"Sending report for restaurant: {restaurant_name}")
    print(f"Phone numbers: {phone_numbers}")

    for phone_number in phone_numbers:
        try:
            # Prepare content variables for the Twilio template
            content_variables = json.dumps({
                "1": restaurant_name,  # Restaurant name
                "2": pdf_blob_filename  # Filename of the PDF (encoded if necessary)
            })

            # Log the content variables for the template
            print(f"Content Variables for WhatsApp: {content_variables}")

            # Send the message using the Twilio template
            message = twilio_client.messages.create(
                to=phone_number,
                content_sid=content_sid,
                content_variables=content_variables,
                messaging_service_sid=messaging_service_sid
            )

            # Log success with message SID
            print(f"Report sent to {phone_number} for {restaurant_name} with SID: {message.sid}")
        except Exception as e:
            # Log any errors during the sending process
            print(f"Failed to send report to {phone_number}: {e}")






# Function to upload PDFs to Azure Blob Storage with a unique name
def upload_pdf_to_blob(restaurant_name, pdf_path):
    # Generate a unique file name by appending the current date
    current_date = datetime.now().strftime("%Y%m%d")  # Format: YYYYMMDD
    blob_name = f"{restaurant_name}_weekly_report_{current_date}.pdf"
    
    # Upload the file to blob storage
    blob_client = blob_service_client.get_blob_client(container=container_name, blob=blob_name)
    with open(pdf_path, "rb") as data:
        blob_client.upload_blob(data, overwrite=True)
    
    print(f"Uploaded {pdf_path} as {blob_name} to Azure Blob Storage.")
    
    # Return the blob URL
    return f"https://daimyoreportstore.blob.core.windows.net/{container_name}/{blob_name}"




# Function to generate insights
def generate_insights(data):
    data = data.copy()  # Create a copy to avoid modifying the original slice
    data['review_questions'] = data['review_questions'].apply(lambda x: json.loads(x.replace("null", "0").replace("'", '"')))
    data['Food_Rating'] = pd.to_numeric(data['review_questions'].apply(lambda x: x.get('Food', 5)), errors='coerce')
    data['Service_Rating'] = pd.to_numeric(data['review_questions'].apply(lambda x: x.get('Service', 5)), errors='coerce')
    data['Atmosphere_Rating'] = pd.to_numeric(data['review_questions'].apply(lambda x: x.get('Atmosphere', 5)), errors='coerce')
    top_authors = data['author_title'].value_counts().head(10)
    rating_distribution = data['review_rating'].value_counts().reindex(range(1, 6), fill_value=0).sort_index()
    insights = {
        'Food': data['Food_Rating'].mean(),
        'Service': data['Service_Rating'].mean(),
        'Atmosphere': data['Atmosphere_Rating'].mean(),
        'Total Reviews': len(data),
        'Average Rating': data['review_rating'].mean(),
        'rating_distribution': rating_distribution,
        'top_authors': top_authors,
        'review_rating': data['review_rating'].mean()
    }
    
    return insights

def print_placeholders(slide):
    """Print available placeholders for debugging."""
    print("Placeholders in this slide:")
    for idx, placeholder in enumerate(slide.placeholders):
        print(f'Placeholder {idx}: {placeholder.placeholder_format.idx}')

def save_plot_to_image(plot_func, filename):
    plt.figure(figsize=(12, 8))
    plot_func()
    plt.tight_layout()
    plt.savefig(filename, format='png')
    plt.close()

def add_textbox(slide, left, top, width, height, text):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear() 
    text_frame.word_wrap = True  # Enable word wrapping
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(16)  # Set font size
    return textbox


def Title(restaurant_name, data, insights):
    slide = prs.slides[0]  # Title Slide layout

    # Access the title shape
    title_shape = slide.shapes.title
    # Set the title text
    title_shape.text = f"{restaurant_name} Review Insights Report"

    # Set the font size to 48 points and make it bold
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(44)
            run.font.name = 'Arial Black'

    # Using the present_week_start and present_week_end for the date range calculation
    date_range_text = f"{present_week_start.strftime('%d %b')} - {present_week_end.strftime('%d %b %Y')}"

    # Adding date range as subtitle
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = date_range_text
    subtitle_shape.text_frame.paragraphs[0].font.name = 'Arial Black'
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)  # Set font size for date range
    subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # Center the text
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(112, 173, 71)

    # Calculate total reviews and average rating for the given time period
    total_reviews = insights['Total Reviews']
    average_rating = insights['Average Rating'] if total_reviews > 0 else 0
    average_rating = round(average_rating, 2)

    # Calculate historic average rating for all time for the given restaurant name
    restaurant_data = data[data['place_name'].str.contains(restaurant_name, case=False, na=False)]
    historic_average = restaurant_data['review_rating'].mean() if len(restaurant_data) > 0 else 0

    # Define the table data
    table_data = [
        ["Total Reviews", f"{total_reviews}"],
        ["Average Rating", f"{average_rating:.2f}"],
        ["Historic Average", f"{historic_average:.2f}"]
    ]

    # Add the table to the slide
    left = Inches(0.75)
    top = Inches(6.5)
    width = Inches(5)
    height = Inches(1.5)
    
    table = slide.shapes.add_table(rows=len(table_data), cols=len(table_data[0]), left=left, top=top, width=width, height=height).table
    # Fill table with data
    for row_idx, row_data in enumerate(table_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            # print(cell)
            cell.text = cell_text
            
            # Apply font size and alignment
            text_frame = cell.text_frame
            p = text_frame.paragraphs[0]
            p.font.size = Pt(20)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # Set cell background colors based on the row
            fill = cell.fill
            fill.solid()
            if row_idx == 0:
                fill.fore_color.rgb = colors["first_row"]
            elif row_idx % 2 == 0:
                fill.fore_color.rgb = colors["even_row"]
            else:
                fill.fore_color.rgb = colors["odd_row"]



def ReviewBreakdown(insights, restaurant_name):
    # Add insights slides
    slide = prs.slides[1]  # Review Breakdown Slide layout
    slide.shapes.title.text = f"Review Breakdown"

    title_shape = slide.shapes.title
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(48)
            run.font.name = 'Arial Black'
    
    # Clear any existing content in the placeholder
    placeholder = slide.placeholders[1]
    placeholder.text_frame.clear()

    # Define the table data
    data = [
        ["Total Reviews", f"{insights['Total Reviews']}"],
        ["Average Customer Rating", f"{insights['Average Rating']:.2f}"],
        ["Food Rating", f"{insights['Food']:.2f}"],
        ["Service Rating", f"{insights['Service']:.2f}"],
        ["Atmosphere Rating", f"{insights['Atmosphere']:.2f}"]
    ]

    # Create a table in the placeholder
    left = Inches(0.55)  # Adjust left position to avoid overlap with title
    top = Inches(3)   # Adjust top position to fit within placeholder
    width = Inches(5.5)   # Adjust width to fit within placeholder
    height = Inches(4) # Adjust height based on content and available space

    table = slide.shapes.add_table(rows=len(data), cols=len(data[0]), left=left, top=top, width=width, height=height).table

    # Set row heights
    for row in range(len(data)):
        table.rows[row].height = Inches(0.3)

    # Fill table with data
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            
            # Apply font size and alignment
            text_frame = cell.text_frame
            p = text_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.bold = True  # Make the font bold
            p.alignment = PP_ALIGN.CENTER

            def is_numeric(s):
                try:
                    float(s)  # Try converting to float
                    return True
                except ValueError:
                    return False
            
            # If the cell contains a rating, make it visually distinct
            if row_idx == 0:
                p.font.color.rgb = RGBColor(255, 255, 255)  # White
                p.font.bold = True  # Bold for headers
                if is_numeric(cell_text):
                    p.font.size = Pt(16)  # Larger font size for numeric cells
                    p.font.bold = True  # Make font bold
            else:
                if is_numeric(cell_text):
                    p.font.size = Pt(16)  # Larger font size for numeric cells
                    p.font.bold = True  # Make font bold
                    p.font.color.rgb = RGBColor(0, 0, 0)  # Set font color (e.g., dark blue)
            fill = cell.fill
            fill.solid()  # Set the fill type to solid for further customization
            
            if row_idx == 0:
                fill.fore_color.rgb = colors["first_row"]
            elif row_idx % 2 == 0:
                fill.fore_color.rgb = colors["even_row"]
            else:
                fill.fore_color.rgb = colors["odd_row"]

    # Add Rating Distribution Graph
    def plot_rating_distribution():
        sns.set_theme(style="whitegrid", context="talk")
        plt.figure(figsize=(10, 6), facecolor='#ffffff')
        ax = sns.barplot(x=insights['rating_distribution'].index, 
                         y=insights['rating_distribution'].values,
                         color='#70AD47',
                         edgecolor='black')
        ax.set_title('Rating Distribution', fontsize=24, weight='bold', color='#000000')
        ax.set_xlabel('Ratings', fontsize=14, color='#000000')
        ax.set_ylabel('Count', fontsize=14, color='#000000')

        for index, value in enumerate(insights['rating_distribution'].values):
            ax.text(index, value + 1, f'{value}', ha='center', fontsize=14, color='black', weight='bold')

        plt.tight_layout()

    rating_dist_img = 'rating_distribution.png'
    save_plot_to_image(plot_rating_distribution, rating_dist_img)

    left_img = Inches(0.5)
    top_img = Inches(5.2)
    width_img = Inches(5)
    height_img = Inches(3.5)
    slide.shapes.add_picture(rating_dist_img, left_img, top_img, width_img, height_img)

    # Add Rating Distribution Insights Text Box
    def generate_rating_distribution_insights(insights, restaurant_name):
        rating_counts = insights['rating_distribution']
        total_reviews = rating_counts.sum()
        most_common_rating = rating_counts.idxmax()
        most_common_rating_count = rating_counts.max()
        high_ratings_percentage = (rating_counts[4] + rating_counts[5]) / total_reviews * 100
        low_ratings_percentage = (rating_counts[1] + rating_counts[2]) / total_reviews * 100
        insights_text=''
        insights_text = (
            f"Rating Distribution Insights for {restaurant_name}:\n\n"
            f"1. {most_common_rating} stars is the most common rating "
            f"({most_common_rating_count} of {total_reviews} reviews).\n\n"
            f"2. High ratings (4-5 stars): {high_ratings_percentage:.1f}% of all reviews.\n\n"
            f"3. Low ratings (1-2 stars): {low_ratings_percentage:.1f}% of all reviews."
        )

        return insights_text

    insights_text = generate_rating_distribution_insights(insights, restaurant_name)

    left_textbox = Inches(1)
    top_textbox = Inches(8.5)
    width_textbox = Inches(5)
    height_textbox = Inches(2)
    add_textbox(slide, left_textbox, top_textbox, width_textbox, height_textbox, insights_text)

def ReviewsTrend(restaurant_data):
    def generate_review_trend_graph(data):
        data['review_date'] = data['review_datetime_utc'].dt.date
        start_date = data['review_date'].min()
        # Group by date and rating
        date_range = pd.date_range(start=start_date, periods=7, freq='D')

        # Group by review_date and review_rating, fill missing dates with 0
        daily_reviews = data.groupby(['review_date', 'review_rating']).size().unstack(fill_value=0)

        # Reindex to ensure all 7 days are present, even if some dates have no reviews
        daily_reviews = daily_reviews.reindex(date_range, fill_value=0)
        # daily_reviews = data.groupby(['review_date', 'review_rating']).size().unstack(fill_value=0)

        # Create the total daily reviews
        daily_reviews['Total'] = daily_reviews.sum(axis=1)
        
        # Calculate average review rating per day
        daily_reviews['Average Rating'] = (daily_reviews.apply(lambda row: sum(r * count for r, count in row.items() if r != 'Total') / row['Total'], axis=1))

        # Create the plot
        fig, ax1 = plt.subplots(figsize=(12, 7))
        
        # Plot total reviews per day as a bar graph with the specified color
        bars = ax1.bar(daily_reviews.index, daily_reviews['Total'], label='Total Reviews', color='#70AD47', edgecolor='black')

        # Add text labels on bars
        for bar in bars:
            height = bar.get_height()
            ax1.text(bar.get_x() + bar.get_width() / 2, height + 0.5, f'{int(height)}', ha='center', va='bottom', fontsize=12, color='black')
        
        # Create a second y-axis for average review rating
        ax2 = ax1.twinx()
        
        # Plot average review rating per day on the second y-axis
        ax2.plot(daily_reviews.index, daily_reviews['Average Rating'], label='Average Rating', color='gold', marker='o', linestyle='-', linewidth=2, markersize=8)
        
        # Add labels for average ratings
        for i in range(len(daily_reviews.index)):
            ax2.text(daily_reviews.index[i], daily_reviews['Average Rating'][i] + 0.1, f"{daily_reviews['Average Rating'][i]:.1f}", 
                    ha='center', va='bottom', fontsize=12, color='black')

        # Set the titles and labels
        ax1.set_title('Daily Review Trend', fontsize=24, weight='bold')
        ax1.set_xlabel('Date', fontsize=14)
        ax1.set_ylabel('Number of Reviews', fontsize=14)
        ax2.set_ylabel('Average Rating', fontsize=14)

        # Format x-axis dates
        ax1.xaxis.set_major_formatter(mdates.DateFormatter('%d\n%b'))  # '%d' for day, '%b' for month abbreviatio
        ax1.xaxis.set_major_locator(mdates.DayLocator())  # Show every day

        # Rotate and format x-axis labels
        # plt.xticks(rotation=45, ha='right', fontsize=3)  # Adjust font size here if needed
        
        # Set y-axis limits and ticks for average rating
        ax2.set_ylim(0, 6)
        ax2.set_yticks(range(0, 7))  # Set ticks from 0 to 6 with a gap of 1
        plt.xticks(rotation=45, ha='right', fontsize=3)
        # Remove gridlines
        ax1.grid(False)
        ax2.grid(False)

        # Add legends for both y-axes
        lines1, labels1 = ax1.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper left', fontsize=8)
        
        # Tight layout to fit all elements
        plt.tight_layout()
        # Find the date with the highest and lowest number of reviews
        most_reviews_date = daily_reviews['Total'].idxmax()
        most_reviews_count = daily_reviews['Total'].max()
        most_reviews_avg_rating = daily_reviews.loc[most_reviews_date, 'Average Rating']

        least_reviews_date = daily_reviews['Total'].idxmin()
        least_reviews_count = daily_reviews['Total'].min()
        least_reviews_avg_rating = daily_reviews.loc[least_reviews_date, 'Average Rating']
        
        insights_text=''
        # Generate dynamic insights text
        insights_text = (
            f"Daily Review Trend Insights for {restaurant_name}:\n\n"
            f"1. The highest number of reviews ({most_reviews_count}) occurred on {most_reviews_date:%d %b}, "
            f"with an average rating of {most_reviews_avg_rating:.1f} stars.\n\n"
            f"2. The lowest number of reviews ({least_reviews_count}) occurred on {least_reviews_date:%d %b}, "
            f"with an average rating of {least_reviews_avg_rating:.1f} stars.\n\n"
            f"3. Overall, the daily average ratings fluctuated, reaching a low of {daily_reviews['Average Rating'].min():.1f} stars "
            f"and peaking at {daily_reviews['Average Rating'].max():.1f} stars."
        )
        return insights_text
    
    # Add Review Trend Slide
    slide = prs.slides[2]  # Review Trend Slide layout
    slide.shapes.title.text = f"Reviews Trend"

    title_shape = slide.shapes.title
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(48)
            run.font.name = 'Arial Black'

    trend_img = 'reviews_trend.png'
    insights_text=generate_review_trend_graph(restaurant_data)
    save_plot_to_image(lambda: generate_review_trend_graph(restaurant_data), trend_img)

    # Define the position and size of the image on the slide
    left = Inches(0.15)
    top = Inches(3.5)
    width = Inches(6)
    height = Inches(4)

    # Add the image to the slide
    slide.shapes.add_picture(trend_img, left, top, width, height)
    # Add a dynamic textbox explaining the graph
    left_textbox = Inches(0.15)  # Align with image left
    top_textbox = Inches(7.75)   # Positioned just below the image
    width_textbox = Inches(6)
    height_textbox = Inches(1)

    textbox = slide.shapes.add_textbox(left_textbox, top_textbox, width_textbox, height_textbox)
    text_frame = textbox.text_frame
    text_frame.clear() 
    text_frame.word_wrap = True

    # Add content to the textbox
    p = text_frame.add_paragraph()
    p.text = insights_text
    p.font.size = Pt(16)

def WeeklyComparison(current_week_data,previous_week_data):
    def weeklies_comparison(current_week_data,previous_week_data):
        current_week_label = f"Current Week\n({present_week_start.strftime('%d %b')} - {present_week_end.strftime('%d %b')})"
        previous_week_label = f"Previous Week\n({previous_week_start.strftime('%d %b')} - {previous_week_end.strftime('%d %b')})"
        # Create a summary dataframe
        summary_data = pd.DataFrame({
            'Week': [previous_week_label, current_week_label],
            'Number of Reviews': [previous_week_data.shape[0], current_week_data.shape[0]],
            'Average Rating': [previous_week_data['review_rating'].mean(), current_week_data['review_rating'].mean()]
        })

        # Plotting
        fig, ax1 = plt.subplots(figsize=(14, 8))  # Increase figure size for clarity

        # Set bar width
        bar_width = 0.4
        x = summary_data['Week']  # X values for bars

        # Plot total reviews per week as a bar graph on the primary Y axis
        bars = ax1.bar(x, summary_data['Number of Reviews'], width=bar_width, label='Number of Reviews', color='#70AD47', edgecolor='black')
        for bar in bars:
            height = bar.get_height()
            ax1.text(bar.get_x() + bar.get_width() / 2, height + 0.5, f'{int(height)}', ha='center', va='bottom', fontsize=12, color='black')
        # Create a secondary Y axis for average rating
        ax2 = ax1.twinx()
        line = ax2.plot(x, summary_data['Average Rating'], label='Average Rating', color='gold', marker='o', linestyle='-', linewidth=2, markersize=8)
        for i, value in enumerate(summary_data['Average Rating']):
            ax2.text(x[i], value + 0.1, f'{value:.2f}', ha='center', va='bottom', fontsize=12, color='black')
        # Labeling
        # ax1.set_xlabel('Week')
        ax1.set_ylabel('Number of Reviews', fontsize=14)
        ax2.set_ylabel('Average Rating', fontsize=14)
        ax2.set_ylim(0, 6)
        ax2.set_yticks(range(0, 7))  # Set ticks from 0 to 6 with a gap of 1
        # Legends
        # Remove gridlines
        ax1.grid(False)
        ax2.grid(False)
        lines1, labels1 = ax1.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper left', fontsize=8)
        # Title
        plt.title('Reviews Summary: Current Week vs Previous Week')
        insight_text=''

        insight_text = f"Weekly Comparison Trend Insights for {restaurant_name}:\n\n"

        # 1. Handle number of reviews (increase or decrease)
        if summary_data['Number of Reviews'][1] > summary_data['Number of Reviews'][0]:
            insight_text += (f"1. Number of reviews increased from {summary_data['Number of Reviews'][0]} "
                            f"to {summary_data['Number of Reviews'][1]}.\n\n")
        elif summary_data['Number of Reviews'][1] < summary_data['Number of Reviews'][0]:
            insight_text += (f"1. Number of reviews decreased from {summary_data['Number of Reviews'][0]} "
                            f"to {summary_data['Number of Reviews'][1]}.\n\n")
        else:
            insight_text += (f"1. Number of reviews remained the same at {summary_data['Number of Reviews'][0]}.\n\n")

        # 2. Handle average rating (increase or decrease)
        if summary_data['Average Rating'][1] > summary_data['Average Rating'][0]:
            insight_text += (f"2. Average rating increased from {summary_data['Average Rating'][0]:.2f} "
                            f"to {summary_data['Average Rating'][1]:.2f}.\n\n")
        elif summary_data['Average Rating'][1] < summary_data['Average Rating'][0]:
            insight_text += (f"2. Average rating decreased from {summary_data['Average Rating'][0]:.2f} "
                            f"to {summary_data['Average Rating'][1]:.2f}.\n\n")
        else:
            insight_text += (f"2. Average rating remained the same at {summary_data['Average Rating'][0]:.2f}.\n\n")

        # 3. General conclusion based on both reviews and rating trends
        if summary_data['Number of Reviews'][1] > summary_data['Number of Reviews'][0] and summary_data['Average Rating'][1] > summary_data['Average Rating'][0]:
            insight_text += "3. Both the number of reviews and the average rating increased, indicating improved customer satisfaction."
        elif summary_data['Number of Reviews'][1] < summary_data['Number of Reviews'][0] and summary_data['Average Rating'][1] < summary_data['Average Rating'][0]:
            insight_text += "3. Both the number of reviews and the average rating decreased, indicating a potential drop in customer satisfaction."
        elif summary_data['Number of Reviews'][1] > summary_data['Number of Reviews'][0] and summary_data['Average Rating'][1] < summary_data['Average Rating'][0]:
            insight_text += "3. The number of reviews increased, but the average rating decreased, suggesting more reviews but with lower satisfaction."
        elif summary_data['Number of Reviews'][1] < summary_data['Number of Reviews'][0] and summary_data['Average Rating'][1] > summary_data['Average Rating'][0]:
            insight_text += "3. The number of reviews decreased, but the average rating increased, indicating fewer reviews but improved satisfaction."
        else:
            insight_text += "3. There is no significant change in both the number of reviews and the average rating."
        
        return insight_text

    # Add the weekly comparison slide
    slide = prs.slides[3]  # Weekly Comparison Slide layout
    slide.shapes.title.text = f"Weekly Comparison"

    title_shape = slide.shapes.title
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(48)
            run.font.name = 'Arial Black'
    
    weekly_trend_img = 'weekly_trend.png'
    insights_text=weeklies_comparison(current_week_data,previous_week_data)
    save_plot_to_image(lambda: weeklies_comparison(current_week_data,previous_week_data), weekly_trend_img)

    # Define the position and size of the image on the slide
    left = Inches(0.25)
    top = Inches(3.5)
    width = Inches(6)
    height = Inches(4)

    # Add the image to the slide
    slide.shapes.add_picture(weekly_trend_img, left, top, width, height)

    # Add a dynamic textbox explaining the graph
    left_textbox = Inches(0.25)  # Align with image left
    top_textbox = Inches(7.75)   # Positioned just below the image
    width_textbox = Inches(6)
    height_textbox = Inches(1)

    textbox = slide.shapes.add_textbox(left_textbox, top_textbox, width_textbox, height_textbox)
    text_frame = textbox.text_frame
    text_frame.clear() 
    text_frame.word_wrap = True

    # Add content to the textbox
    p = text_frame.add_paragraph()
    p.text = insights_text
    p.font.size = Pt(16)

# Function to generate summaries of positive and negative reviews
def generate_review_summaries(restaurant_data):
    # Filter positive reviews (4 & 5 stars)
    positive_reviews = restaurant_data[restaurant_data['review_rating'].isin([4, 5])]['review_text'].dropna()
    positive_text = " ".join(positive_reviews.tolist())

    # Filter negative reviews (1, 2 & 3 stars)
    negative_reviews = restaurant_data[restaurant_data['review_rating'].isin([1, 2, 3])]['review_text'].dropna()
    negative_text = " ".join(negative_reviews.tolist())
    positive_summary=""
    negative_summary=""
    # Generate summary for positive reviews
    try:
        positive_summary = summary_chain.invoke({"Reviews": positive_text})
    except Exception as e:
        print(f"Failed to generate summary for positive reviews: {e}")
        positive_summary = "Could not generate summary for positive reviews."

    # Generate summary for negative reviews
    try:
        negative_summary = summary_chain.invoke({"Reviews": negative_text})
    except Exception as e:
        print(f"Failed to generate summary for negative reviews: {e}")
        negative_summary = "Could not generate summary for negative reviews."

    return positive_summary, negative_summary

def add_textbox(slide, left, top, width, height, text, font_size=16):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear() 
    text_frame.word_wrap = True  # Enable word wrapping
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)  # Set font size
    p.alignment = PP_ALIGN.LEFT  # Align text to the left
    return textbox

def SentimentAnalysisPositive(restaurant_data):
    # Generate summary of positive reviews
    positive_summary, _ = generate_review_summaries(restaurant_data)
    
    # Use the existing slide layout for positive sentiment analysis (assuming it is the 6th slide in the template)
    slide = prs.slides[4]  # Positive Sentiment Analysis Slide layout
    slide.shapes.title.text = "Positive Reviews Sentiment Analysis"

    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(40)
            run.font.name = 'Arial Black'
            run.font.bold = True

    # Adjust the position and size to fit within the 31 cm height and 16 cm width
    left = Inches(0.5)  # Adjust the left margin if needed
    top = Inches(3.0)     # Top margin
    width = Inches(5.5)   # Adjusted width to prevent overflow
    height = Inches(7)  # Increased height for more text space
    
    # Add the text box for the positive reviews summary
    add_textbox(slide, left, top, width, height, f"Positive Reviews Summary: {positive_summary}", font_size=16)

def SentimentAnalysisNegative(restaurant_data):
    # Generate summary of negative reviews
    _, negative_summary = generate_review_summaries(restaurant_data)
    
    # Use the existing slide layout for negative sentiment analysis (assuming it is the 7th slide in the template)
    slide = prs.slides[5]  # Negative Sentiment Analysis Slide layout
    slide.shapes.title.text = "Negative Reviews Sentiment Analysis"

    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(40)
            run.font.name = 'Arial Black'
            run.font.bold = True

    # Adjust the position and size to fit within the 31 cm height and 16 cm width
    left = Inches(0.5)  # Adjust the left margin if needed
    top = Inches(3.0)     # Top margin
    width = Inches(5.5)   # Adjusted width to prevent overflow
    height = Inches(7)  # Increased height for more text space
    
    # Add the text box for the negative reviews summary
    add_textbox(slide, left, top, width, height, f"Negative Reviews Summary: {negative_summary}", font_size=16)

def ThankYou():
    slide = prs.slides[6]  # Adjusted to the 8th slide index (7th in 0-based index)

    # Access the title shape
    title_shape = slide.shapes.title
    title_shape.text = "THANK YOU!"  # Set text to all caps

    # Center the title text and set the font style
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(60)  # Increase font size for prominence
            run.font.name = 'Arial Black'
            run.font.bold = True
            run.font.color.rgb = RGBColor(112, 173, 71)  # Green Accent 6, Darker 25% color

    # Add a paragraph for the additional message and center it
    text_frame = title_shape.text_frame
    text_frame.add_paragraph()  # Add a blank paragraph to create some space
    p = text_frame.add_paragraph()
    p.text = "Crafted with care by Daimyo"
    p.alignment = PP_ALIGN.CENTER  # Center the additional text
    p.space_before = Pt(80)  # Add some space before the paragraph
    for run in p.runs:
        run.font.size = Pt(30)  # Font size for the additional message
        run.font.name = 'Arial Black'
        run.font.bold = True

# Function to generate restaurant presentation (including insights)
def generate_restaurant_presentation(restaurant_name, data, current_week_data, previous_week_data):
    insights = generate_insights(current_week_data)
    # prs = Presentation("TemplatePlain.pptx")  # Make sure your template has the correct dimensions

    # 1. TITLE SLIDE
    Title(restaurant_name, data,insights)

    # 2. Review Breakdown SLIDE
    ReviewBreakdown(insights, restaurant_name)

    # 3. Review Trend SLIDE
    ReviewsTrend(current_week_data)  # Pass data for review trend

    # 4. Weekly Comparison SLIDE
    WeeklyComparison(current_week_data,previous_week_data)  # Pass data for weekly comparison

    # 5. Positive Sentiment Analysis SLIDE
    SentimentAnalysisPositive(current_week_data)  # Pass data for positive sentiment analysis

    # 6. Negative Sentiment Analysis SLIDE
    SentimentAnalysisNegative(current_week_data)  # Pass data for negative sentiment analysis

    # 7. Thank You SLIDE
    ThankYou()

    # Save the presentation as a PDF
    ppt_path = f'{restaurant_name}_review_presentation.pptx'
    prs.save(ppt_path)
    print(f"Presentation for {restaurant_name} generated successfully!")

    # Convert to PDF
    pdf_path = f'{restaurant_name}_review.pdf'
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
    presentation.SaveAs(os.path.abspath(pdf_path), FileFormat=32)
    presentation.Close()
    powerpoint.Quit()
    print(f"Presentation saved as PDF: {pdf_path}")
    
    return pdf_path

# Fetch data from the database
data = fetch_data_from_db()

# Set up the timezone to Indian Standard Time
custom_offset = pytz.FixedOffset(330)  # 5:30 offset
data['sentiment_score'] = data['sentiment_score'].fillna(0).astype(float)
data['review_datetime_utc'] = pd.to_datetime(data['review_datetime_utc'], errors='coerce')
data['review_datetime_utc'] = data['review_datetime_utc'].dt.tz_localize('UTC').dt.tz_convert(custom_offset)

# Report type (for weekly reports)
report_type = "customer_review_report_weekly"

# Generate reports for all restaurants
restaurant_names = fetch_restro_from_db()
for restaurant_name in restaurant_names:

    prs = Presentation("TemplatePlain.pptx")  # Make sure your template has the correct dimensions
    restaurant_data = filter_by_restaurant(data, restaurant_name)
    # Check if a report has already been generated this week
    # if check_report_generated(restaurant_name, report_type):
    #     print(f"Reports for {restaurant_name} already sent for this week.")
    #     continue

    # Filter data by restaurant
    restaurant_data = data[data['place_name'].str.contains(restaurant_name, case=False, na=False)]
    if restaurant_data.empty:
        print(f"No data found for {restaurant_name}.")
        continue

    # Date calculations (unchanged)
    now = datetime.now(custom_offset).date()
    present_week_end = now - timedelta(days=1)
    present_week_start = present_week_end - timedelta(days=6)
    previous_week_end = present_week_start - timedelta(days=1)
    previous_week_start = previous_week_end - timedelta(days=6)
    present_week_start = datetime.combine(present_week_start, datetime.min.time(), tzinfo=custom_offset)
    present_week_end = datetime.combine(present_week_end, datetime.max.time(), tzinfo=custom_offset)
    previous_week_start = datetime.combine(previous_week_start, datetime.min.time(), tzinfo=custom_offset)
    previous_week_end = datetime.combine(previous_week_end, datetime.max.time(), tzinfo=custom_offset)

    current_week_data = restaurant_data[(restaurant_data['review_datetime_utc'] >= present_week_start) & (restaurant_data['review_datetime_utc'] <= present_week_end)]
    previous_week_data = restaurant_data[(restaurant_data['review_datetime_utc'] >= previous_week_start) & (restaurant_data['review_datetime_utc'] <= previous_week_end)]

    # Generate the presentation and save as PDF
    pdf_path = generate_restaurant_presentation(restaurant_name, restaurant_data, current_week_data, previous_week_data)

    # Upload the PDF to Azure Blob Storage and get the blob URL
    pdf_blob_url = upload_pdf_to_blob(restaurant_name, pdf_path)

    # Send the report via WhatsApp
    send_report_via_whatsapp(restaurant_name, pdf_blob_url)

    # Update the report tracking
    update_report_generated_date(restaurant_name, report_type)
